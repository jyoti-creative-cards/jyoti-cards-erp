#!/usr/bin/env python3
"""Generate Synvoric capability deck (PPT) and legal document templates."""

from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
TEMPLATES = ASSETS / "templates"
LOGO = ROOT / "public" / "logo.png"

NAVY = RGBColor(0x05, 0x1C, 0x2C)
BLUE = RGBColor(0x00, 0x66, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x33, 0x41, 0x55)

SERVICES = [
    ("Frontend Developer Staffing", "React, Angular, Vue, Next.js specialists for production-ready UI."),
    ("AI & ML Engineer Staffing", "ML engineers, LLM integrators, and data scientists."),
    ("Web App Development", "MERN, MEAN, and enterprise SaaS teams."),
    ("Mobile App Development", "iOS, Android, React Native, and Flutter teams."),
    ("Dedicated Product Teams", "Managed squads with tech lead, devs, and QA."),
    ("End-to-End Software Delivery", "Discovery through deployment and maintenance."),
]

WHY = [
    ("Rigorous Vetting", "Technical assessments, live coding, and communication screening."),
    ("Speed to Hire", "Average two weeks from requirement to onboarded developer."),
    ("Flexible Models", "Staff augmentation, dedicated teams, or full delivery."),
    ("Account Partnership", "Dedicated account manager with weekly check-ins."),
]

STATS = [("50+", "Developers Placed"), ("10+", "Clients"), ("100%", "Retention"), ("2 Weeks", "Time to Hire")]

TECH = {
    "Frontend": "React · Next.js · Angular · Vue.js · TypeScript · Tailwind",
    "AI & ML": "Python · TensorFlow · PyTorch · OpenAI · Hugging Face · MLOps",
    "Mobile": "React Native · Flutter · Swift · Kotlin · iOS · Android",
    "Backend": "Node.js · MongoDB · PostgreSQL · GraphQL · AWS · Docker",
}

CLIENTS = [
    "Chandrawat Concrete Works", "Firangi Cafe & Bar", "The Grand Shaurya",
    "Techturtle", "SIG SIGMA", "Tripolic", "Jyoti Creative Cards",
]


def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, num: int, total: int, dark: bool = False):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.35), Inches(10), Inches(0.25)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY if not dark else RGBColor(0x03, 0x14, 0x22)
    bar.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.5), Inches(5.38), Inches(4), Inches(0.2))
    tf = left.text_frame
    p = tf.paragraphs[0]
    p.text = "Synvoric · Confidential"
    p.font.size = Pt(8)
    p.font.color.rgb = SLATE if not dark else RGBColor(0x94, 0xA3, 0xB8)

    right = slide.shapes.add_textbox(Inches(8.5), Inches(5.38), Inches(1), Inches(0.2))
    rp = right.text_frame.paragraphs[0]
    rp.text = f"{num:02d} / {total:02d}"
    rp.font.size = Pt(8)
    rp.font.bold = True
    rp.font.color.rgb = BLUE if not dark else RGBColor(0x66, 0x99, 0xFF)
    rp.alignment = PP_ALIGN.RIGHT


def add_title_block(slide, eyebrow: str, title: str):
    eb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(9), Inches(0.3))
    ep = eb.text_frame.paragraphs[0]
    ep.text = eyebrow.upper()
    ep.font.size = Pt(9)
    ep.font.bold = True
    ep.font.color.rgb = BLUE
    ep.font.name = "Calibri"

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(9), Inches(0.6))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    tp.font.name = "Calibri"

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.35), Inches(0.45), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9
    blank = prs.slide_layouts[6]
    total = 12

    # 01 Cover
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(2.8), Inches(1.2), width=Inches(4.4))
    tag = s.shapes.add_textbox(Inches(1), Inches(2.85), Inches(8), Inches(0.35))
    tp = tag.text_frame.paragraphs[0]
    tp.text = "AI · DATA · ENGINEERING · SOLUTIONS"
    tp.font.size = Pt(11)
    tp.font.color.rgb = BLUE
    tp.font.bold = True
    tp.alignment = PP_ALIGN.CENTER
    sub = s.shapes.add_textbox(Inches(1), Inches(3.35), Inches(8), Inches(0.5))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "Capability Overview"
    sp.font.size = Pt(22)
    sp.font.color.rgb = WHITE
    sp.alignment = PP_ALIGN.CENTER
    dt = s.shapes.add_textbox(Inches(1), Inches(3.85), Inches(8), Inches(0.3))
    dp = dt.text_frame.paragraphs[0]
    dp.text = f"May {date.today().year}  ·  synvoric.com"
    dp.font.size = Pt(10)
    dp.font.color.rgb = SLATE
    dp.alignment = PP_ALIGN.CENTER
    add_footer(s, 1, total, dark=True)

    # 02 About
    s = prs.slides.add_slide(blank)
    add_title_block(s, "About Us", "Who We Are")
    body = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(5.2), Inches(3.5))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (
        "Synvoric is a professional IT staffing company specializing in frontend "
        "developers and AI/ML engineers. We partner with enterprises and growth-stage "
        "companies to deliver vetted engineering talent for web applications, mobile "
        "apps, and AI-powered products."
    )
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = (
        "\nFounded on the belief that great software starts with great people, we "
        "combine rigorous technical vetting with a partnership-first approach."
    )
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK
    for i, (val, lbl) in enumerate(STATS):
        col, row = i % 2, i // 2
        x = Inches(6.2 + col * 1.85)
        y = Inches(1.7 + row * 1.35)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.65), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        btf = box.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = btf.paragraphs[0]
        vp.text = val
        vp.font.size = Pt(22)
        vp.font.bold = True
        vp.font.color.rgb = BLUE
        vp.alignment = PP_ALIGN.CENTER
        lp = btf.add_paragraph()
        lp.text = lbl.upper()
        lp.font.size = Pt(8)
        lp.font.color.rgb = SLATE
        lp.alignment = PP_ALIGN.CENTER
    add_footer(s, 2, total)

    # 03 Divider
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    num = s.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(3), Inches(0.8))
    num.text_frame.paragraphs[0].text = "02"
    num.text_frame.paragraphs[0].font.size = Pt(60)
    num.text_frame.paragraphs[0].font.bold = True
    num.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    num.text_frame.paragraphs[0].font.color.brightness = 0.15
    tit = s.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(8), Inches(0.6))
    tit.text_frame.paragraphs[0].text = "Our Services"
    tit.text_frame.paragraphs[0].font.size = Pt(32)
    tit.text_frame.paragraphs[0].font.bold = True
    tit.text_frame.paragraphs[0].font.color.rgb = WHITE
    add_footer(s, 3, total, dark=True)

    # 04 Services grid
    s = prs.slides.add_slide(blank)
    add_title_block(s, "Services", "Staffing & Development Solutions")
    for i, (title, desc) in enumerate(SERVICES):
        col, row = i % 3, i // 3
        x, y = Inches(0.55 + col * 3.1), Inches(1.55 + row * 1.85)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.9), Inches(1.65))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.9), Inches(0.04))
        top.fill.solid()
        top.fill.fore_color.rgb = BLUE
        top.line.fill.background()
        btf = box.text_frame
        btf.margin_left = DocInches(0.12)
        btf.margin_top = DocInches(0.12)
        btf.word_wrap = True
        hp = btf.paragraphs[0]
        hp.text = title
        hp.font.size = Pt(11)
        hp.font.bold = True
        hp.font.color.rgb = NAVY
        dp = btf.add_paragraph()
        dp.text = desc
        dp.font.size = Pt(9)
        dp.font.color.rgb = DARK
    add_footer(s, 4, total)

    # 05-06 Service details (frontend + AI)
    details = [
        ("Frontend Developer Staffing", [
            "Responsive web & SPA development",
            "Design system & component architecture",
            "Performance optimization & accessibility",
            "Senior React, Next.js, Angular developers",
        ]),
        ("AI & ML Engineer Staffing", [
            "LLM integration & prompt engineering",
            "Custom model training & fine-tuning",
            "MLOps & deployment pipelines",
            "NLP, computer vision & analytics",
        ]),
    ]
    for idx, (title, bullets) in enumerate(details):
        s = prs.slides.add_slide(blank)
        add_title_block(s, "Service Detail", title)
        bx = s.shapes.add_textbox(Inches(0.6), Inches(1.55), Inches(8.8), Inches(3.2))
        btf = bx.text_frame
        for j, b in enumerate(bullets):
            p = btf.paragraphs[0] if j == 0 else btf.add_paragraph()
            p.text = f"—  {b}"
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.space_after = Pt(10)
        add_footer(s, 5 + idx, total)

    # 07 Web & Mobile
    s = prs.slides.add_slide(blank)
    add_title_block(s, "Service Detail", "Web & Mobile App Development")
    for i, (title, bullets) in enumerate([
        ("Web App Development", ["MERN & MEAN stack", "Enterprise portals & SaaS", "REST & GraphQL APIs"]),
        ("Mobile App Development", ["iOS (Swift) & Android (Kotlin)", "React Native & Flutter", "App Store deployment"]),
    ]):
        x = Inches(0.6 + i * 4.7)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.55), Inches(4.4), Inches(3.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        btf = box.text_frame
        btf.margin_left = DocInches(0.15)
        btf.margin_top = DocInches(0.15)
        hp = btf.paragraphs[0]
        hp.text = title
        hp.font.bold = True
        hp.font.size = Pt(13)
        hp.font.color.rgb = NAVY
        for b in bullets:
            bp = btf.add_paragraph()
            bp.text = f"—  {b}"
            bp.font.size = Pt(11)
            bp.font.color.rgb = DARK
    add_footer(s, 7, total)

    # 08 Divider Tech
    s = prs.slides.add_slide(blank)
    set_slide_bg(s, NAVY)
    num = s.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(3), Inches(0.8))
    num.text_frame.paragraphs[0].text = "03"
    num.text_frame.paragraphs[0].font.size = Pt(60)
    num.text_frame.paragraphs[0].font.bold = True
    num.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x30, 0x40, 0x55)
    tit = s.shapes.add_textbox(Inches(0.6), Inches(3.85), Inches(8), Inches(0.6))
    tit.text_frame.paragraphs[0].text = "Technologies"
    tit.text_frame.paragraphs[0].font.size = Pt(32)
    tit.text_frame.paragraphs[0].font.bold = True
    tit.text_frame.paragraphs[0].font.color.rgb = WHITE
    add_footer(s, 8, total, dark=True)

    # 09 Tech
    s = prs.slides.add_slide(blank)
    add_title_block(s, "Tech Stack", "Technologies Our Developers Master")
    for i, (cat, items) in enumerate(TECH.items()):
        x = Inches(0.55 + i * 2.35)
        box = s.shapes.add_textbox(x, Inches(1.55), Inches(2.2), Inches(3.5))
        btf = box.text_frame
        btf.word_wrap = True
        cp = btf.paragraphs[0]
        cp.text = cat.upper()
        cp.font.size = Pt(9)
        cp.font.bold = True
        cp.font.color.rgb = BLUE
        ip = btf.add_paragraph()
        ip.text = items
        ip.font.size = Pt(10)
        ip.font.color.rgb = DARK
    add_footer(s, 9, total)

    # 10 Why Synvoric
    s = prs.slides.add_slide(blank)
    add_title_block(s, "Differentiators", "Why Synvoric")
    for i, (title, body) in enumerate(WHY):
        col, row = i % 2, i // 2
        x = Inches(0.6 + col * 4.7)
        y = Inches(1.55 + row * 1.85)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.05), Inches(1.55))
        bar.fill.solid()
        bar.fill.fore_color.rgb = BLUE
        bar.line.fill.background()
        bx = s.shapes.add_textbox(x + Inches(0.15), y, Inches(4.3), Inches(1.55))
        btf = bx.text_frame
        btf.word_wrap = True
        hp = btf.paragraphs[0]
        hp.text = title
        hp.font.bold = True
        hp.font.size = Pt(12)
        hp.font.color.rgb = NAVY
        bp = btf.add_paragraph()
        bp.text = body
        bp.font.size = Pt(10)
        bp.font.color.rgb = DARK
    add_footer(s, 10, total)

    # 11 Process + Clients
    s = prs.slides.add_slide(blank)
    add_title_block(s, "Track Record", "How We Work & Who We Serve")
    steps = [("01", "Discovery"), ("02", "Match"), ("03", "Onboard"), ("04", "Support")]
    for i, (st, lbl) in enumerate(steps):
        x = Inches(0.6 + i * 2.35)
        bx = s.shapes.add_textbox(x, Inches(1.55), Inches(2.1), Inches(1.2))
        btf = bx.text_frame
        np = btf.paragraphs[0]
        np.text = st
        np.font.size = Pt(24)
        np.font.bold = True
        np.font.color.rgb = BLUE
        np.alignment = PP_ALIGN.CENTER
        lp = btf.add_paragraph()
        lp.text = lbl
        lp.font.size = Pt(11)
        lp.font.bold = True
        lp.font.color.rgb = NAVY
        lp.alignment = PP_ALIGN.CENTER
    cl = s.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(9), Inches(1.8))
    ctf = cl.text_frame
    ctf.word_wrap = True
    cp = ctf.paragraphs[0]
    cp.text = "SELECTED CLIENTS"
    cp.font.size = Pt(9)
    cp.font.bold = True
    cp.font.color.rgb = BLUE
    ip = ctf.add_paragraph()
    ip.text = "  ·  ".join(CLIENTS)
    ip.font.size = Pt(10)
    ip.font.color.rgb = DARK
    add_footer(s, 11, total)

    # 12 Contact
    s = prs.slides.add_slide(blank)
    add_title_block(s, "Get in Touch", "Let's Build Your Engineering Team")
    cx = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(2.5))
    ctf = cx.text_frame
    ep = ctf.paragraphs[0]
    ep.text = "anshul.agrawal@synvoric.com"
    ep.font.size = Pt(20)
    ep.font.bold = True
    ep.font.color.rgb = BLUE
    ep.alignment = PP_ALIGN.CENTER
    wp = ctf.add_paragraph()
    wp.text = "www.synvoric.com"
    wp.font.size = Pt(14)
    wp.font.color.rgb = SLATE
    wp.alignment = PP_ALIGN.CENTER
    np = ctf.add_paragraph()
    np.text = "\nStaff Augmentation  ·  Dedicated Teams  ·  Project Delivery"
    np.font.size = Pt(11)
    np.font.color.rgb = DARK
    np.alignment = PP_ALIGN.CENTER
    add_footer(s, 12, total)

    out = ASSETS / "Synvoric-Capability-Deck.pptx"
    prs.save(str(out))
    print(f"Created {out}")


# ── Word document helpers ──

def doc_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = DocRGB(0x05, 0x1C, 0x2C)
    return h


def doc_para(doc, text, bold=False, space_after=6):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size = DocPt(11)
    run.font.name = "Calibri"
    run.bold = bold
    p.paragraph_format.space_after = DocPt(space_after)
    return p


def doc_bullet(doc, text):
    p = doc.add_paragraph(text, style="List Bullet")
    for run in p.runs:
        run.font.size = DocPt(11)
        run.font.name = "Calibri"
    return p


def doc_disclaimer(doc):
    p = doc.add_paragraph()
    run = p.add_run(
        "DISCLAIMER: This document is a template only and does not constitute legal advice. "
        "Synvoric recommends review by qualified legal counsel before execution."
    )
    run.font.size = DocPt(9)
    run.font.italic = True
    run.font.color.rgb = DocRGB(0x64, 0x74, 0x8B)
    p.paragraph_format.space_after = DocPt(12)


def build_proposal():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = DocInches(1)
    section.bottom_margin = DocInches(1)
    section.left_margin = DocInches(1)
    section.right_margin = DocInches(1)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("SYNVORIC")
    r.font.size = DocPt(22)
    r.font.bold = True
    r.font.color.rgb = DocRGB(0x05, 0x1C, 0x2C)
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = sub.add_run("Client Proposal")
    sr.font.size = DocPt(16)
    sr.font.color.rgb = DocRGB(0x00, 0x66, 0xFF)

    doc_disclaimer(doc)
    doc_para(doc, "Date: [PROPOSAL DATE]")
    doc_para(doc, "Prepared for: [CLIENT COMPANY NAME]")
    doc_para(doc, "Attention: [CLIENT CONTACT NAME], [TITLE]")
    doc_para(doc, "Prepared by: Synvoric  ·  anshul.agrawal@synvoric.com  ·  www.synvoric.com")
    doc.add_paragraph()

    doc_heading(doc, "1. Executive Summary", 2)
    doc_para(doc, (
        "Synvoric proposes to provide [ENGAGEMENT TYPE: Staff Augmentation / Dedicated Team / "
        "Project Delivery] services to [CLIENT COMPANY NAME] for [PROJECT OR TEAM DESCRIPTION]. "
        "This proposal outlines our understanding of your requirements, proposed approach, "
        "team composition, timeline, and commercial terms."
    ))

    doc_heading(doc, "2. Understanding of Requirements", 2)
    doc_para(doc, "[CLIENT COMPANY NAME] requires the following:")
    for item in [
        "[Describe business objective and context]",
        "[Required roles — e.g., 2 Senior React Developers, 1 ML Engineer]",
        "[Technology stack — e.g., React, Node.js, Python, AWS]",
        "[Engagement duration — e.g., 6 months, renewable]",
        "[Working model — remote / hybrid / on-site]",
    ]:
        doc_bullet(doc, item)

    doc_heading(doc, "3. Proposed Solution", 2)
    doc_para(doc, "Synvoric will deliver the following:")
    doc_bullet(doc, "Pre-vetted developers matched to your stack and culture")
    doc_bullet(doc, "Onboarding within [X] business days of agreement signing")
    doc_bullet(doc, "Dedicated account manager for ongoing coordination")
    doc_bullet(doc, "Replacement guarantee if a resource does not meet expectations within [30] days")

    doc_heading(doc, "4. Proposed Team", 2)
    table = doc.add_table(rows=4, cols=4)
    table.style = "Table Grid"
    headers = ["Role", "Experience", "Skills", "Allocation"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    rows_data = [
        ["[Role 1]", "[X+ years]", "[Skills]", "[Full-time / Part-time]"],
        ["[Role 2]", "[X+ years]", "[Skills]", "[Full-time / Part-time]"],
        ["[Role 3]", "[X+ years]", "[Skills]", "[Full-time / Part-time]"],
    ]
    for ri, row in enumerate(rows_data, 1):
        for ci, val in enumerate(row):
            table.rows[ri].cells[ci].text = val

    doc.add_paragraph()
    doc_heading(doc, "5. Timeline & Milestones", 2)
    for item in [
        "Week 1: Agreement execution and requirement finalization",
        "Week 2: Candidate shortlisting and client interviews",
        "Week 3: Selected developers onboarded and integrated",
        "Ongoing: Weekly status reports and performance reviews",
    ]:
        doc_bullet(doc, item)

    doc_heading(doc, "6. Commercial Terms", 2)
    doc_para(doc, "The following commercial terms are proposed (subject to MSA execution):")
    doc_bullet(doc, "Billing Rate: [INR/USD] [AMOUNT] per [developer/month or hour]")
    doc_bullet(doc, "Billing Cycle: Monthly in arrears, invoice due within [15/30] days")
    doc_bullet(doc, "Engagement Term: [X months], auto-renewable unless terminated with [30] days notice")
    doc_bullet(doc, "Expenses: [Included / Reimbursable as pre-approved]")
    doc_bullet(doc, "Payment Method: Bank transfer to Synvoric designated account")

    doc_heading(doc, "7. Why Synvoric", 2)
    for title, body in WHY:
        doc_para(doc, f"{title}: {body}")

    doc_heading(doc, "8. Terms & Acceptance", 2)
    doc_para(doc, (
        "This proposal is valid for [30] days from the date above. Services shall be governed "
        "by the Synvoric Master Service Agreement (MSA) upon acceptance. To proceed, please "
        "sign below and return to anshul.agrawal@synvoric.com."
    ))
    doc.add_paragraph()
    doc_para(doc, "Accepted by [CLIENT COMPANY NAME]:")
    doc_para(doc, "Name: ___________________________    Title: ___________________________")
    doc_para(doc, "Signature: _______________________    Date: ___________________________")

    out = TEMPLATES / "Synvoric-Client-Proposal-Template.docx"
    doc.save(str(out))
    print(f"Created {out}")


def build_msa():
    doc = Document()
    doc_disclaimer(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("MASTER SERVICE AGREEMENT")
    r.font.size = DocPt(18)
    r.font.bold = True
    r.font.color.rgb = DocRGB(0x05, 0x1C, 0x2C)

    doc_para(doc, "Effective Date: [EFFECTIVE DATE]")
    doc_para(doc, 'This Master Service Agreement ("Agreement") is entered into between:')
    doc_para(doc, 'Provider: Synvoric ("Synvoric", "Provider", "we")', bold=True)
    doc_para(doc, "Contact: anshul.agrawal@synvoric.com  ·  www.synvoric.com")
    doc_para(doc, 'Client: [CLIENT COMPANY NAME] ("Client", "you")', bold=True)
    doc_para(doc, "Address: [CLIENT ADDRESS]")
    doc_para(doc, 'Synvoric and Client are each a "Party" and collectively the "Parties."')

    sections = [
        ("1. Services", [
            'Synvoric shall provide IT staffing, staff augmentation, dedicated development teams, and/or software development services ("Services") as described in one or more Statements of Work ("SOW") or work orders executed under this Agreement.',
            "Each SOW shall specify scope, deliverables, team composition, rates, and timeline. In case of conflict, the SOW prevails over this Agreement for that engagement.",
        ]),
        ("2. Term", [
            "This Agreement commences on the Effective Date and continues until terminated in accordance with Section 12.",
            "Individual SOWs may specify their own term within the life of this Agreement.",
        ]),
        ("3. Fees & Payment", [
            "Client shall pay fees as set forth in the applicable SOW. Unless otherwise stated, fees are billed monthly in arrears.",
            "Invoices are due within [15/30] days of receipt. Late payments accrue interest at [1.5]% per month or the maximum permitted by law.",
            "Client shall reimburse pre-approved out-of-pocket expenses with supporting documentation.",
            "All fees are exclusive of applicable taxes, duties, and withholdings, which Client shall bear unless Synvoric is required by law to collect them.",
            "Synvoric may suspend Services for overdue amounts exceeding [30] days after written notice.",
        ]),
        ("4. Client Responsibilities", [
            "Provide timely access to systems, documentation, and personnel necessary for Service delivery.",
            "Designate a primary point of contact for day-to-day coordination.",
            "Review and approve deliverables within agreed timeframes.",
            "Comply with all applicable laws in use of Services and Synvoric personnel.",
        ]),
        ("5. Personnel", [
            "Synvoric personnel (including contractors) remain employees or contractors of Synvoric unless otherwise agreed in writing.",
            "Client shall not directly hire or engage Synvoric personnel outside this Agreement during the term and for [12] months following termination (see Section 10 — Non-Solicitation).",
            "Synvoric shall use reasonable efforts to replace personnel who are unavailable or underperforming.",
        ]),
        ("6. Confidentiality", [
            '"Confidential Information" means non-public business, technical, financial, or personnel information disclosed by either Party.',
            "Each Party shall: (a) use Confidential Information only for purposes of this Agreement; (b) protect it with at least reasonable care; (c) not disclose it to third parties except to personnel and advisors with a need to know who are bound by similar obligations.",
            "Exclusions: information that is public, independently developed, rightfully received without restriction, or required to be disclosed by law (with prompt notice where permitted).",
            "Confidentiality obligations survive termination for [3] years, except trade secrets which survive as long as protected by law.",
        ]),
        ("7. Intellectual Property", [
            "Pre-Existing IP: Each Party retains ownership of IP it owned or developed prior to or outside the Services.",
            'Work Product: Unless the SOW states otherwise, all deliverables, code, documentation, and materials created specifically for Client under a paid SOW ("Work Product") shall be owned by Client upon full payment.',
            "Synvoric retains ownership of its general methodologies, tools, templates, and know-how used in delivering Services.",
            "Synvoric grants Client a perpetual, royalty-free license to use Synvoric pre-existing IP embedded in Work Product to the extent necessary to use the Work Product.",
            "Client grants Synvoric a limited license to use Client materials solely to perform Services.",
        ]),
        ("8. Warranties & Disclaimer", [
            "Synvoric warrants Services will be performed in a professional manner consistent with industry standards.",
            "EXCEPT AS EXPRESSLY SET FORTH HEREIN, SYNVORIC DISCLAIMS ALL OTHER WARRANTIES, EXPRESS OR IMPLIED, INCLUDING MERCHANTABILITY AND FITNESS FOR A PARTICULAR PURPOSE.",
            "Synvoric does not warrant uninterrupted or error-free Services.",
        ]),
        ("9. Limitation of Liability", [
            "NEITHER PARTY SHALL BE LIABLE FOR INDIRECT, INCIDENTAL, SPECIAL, CONSEQUENTIAL, OR PUNITIVE DAMAGES.",
            "Each Party's aggregate liability under this Agreement shall not exceed the fees paid or payable by Client in the [12] months preceding the claim, except for breaches of confidentiality, IP infringement, or willful misconduct.",
        ]),
        ("10. Non-Solicitation", [
            "During the term and for [12] months after termination, Client shall not directly or indirectly solicit, hire, or engage any Synvoric personnel who performed Services under this Agreement without Synvoric's prior written consent.",
            "If Client hires such personnel in breach of this section, Client shall pay Synvoric a placement fee equal to [12] months of the individual's gross compensation with Synvoric, or [INR/USD AMOUNT], whichever is greater.",
            "General public job postings not targeted at Synvoric personnel are excluded.",
        ]),
        ("11. Indemnification", [
            "Synvoric shall indemnify Client against third-party claims that Work Product created solely by Synvoric infringes IP rights, subject to prompt notice and cooperation.",
            "Client shall indemnify Synvoric against claims arising from Client materials, instructions, or unlawful use of Services.",
        ]),
        ("12. Termination", [
            "Either Party may terminate this Agreement or any SOW with [30] days written notice.",
            "Either Party may terminate immediately for material breach not cured within [15] days of written notice.",
            "Upon termination, Client shall pay for Services rendered through the termination date. Sections intended to survive (Confidentiality, IP, Non-Solicitation, Limitation of Liability) shall survive.",
        ]),
        ("13. Governing Law & Disputes", [
            "This Agreement is governed by the laws of [India / State of ____________], without regard to conflict of law principles.",
            "Disputes shall first be resolved through good-faith negotiation. Failing resolution within [30] days, disputes shall be subject to the exclusive jurisdiction of courts in [City, State/Country].",
        ]),
        ("14. General", [
            "This Agreement constitutes the entire agreement and supersedes prior discussions.",
            "Amendments must be in writing signed by both Parties.",
            "Neither Party may assign without consent, except to an affiliate or successor.",
            "Notices shall be sent to the addresses above (or email for Synvoric: anshul.agrawal@synvoric.com).",
        ]),
    ]

    for heading, paragraphs in sections:
        doc_heading(doc, heading, 2)
        for para in paragraphs:
            doc_para(doc, para)

    doc.add_paragraph()
    doc_heading(doc, "Signatures", 2)
    doc_para(doc, "IN WITNESS WHEREOF, the Parties have executed this Agreement as of the Effective Date.")
    doc.add_paragraph()
    doc_para(doc, "SYNVORIC                              CLIENT")
    doc_para(doc, "Name: _______________________          Name: _______________________")
    doc_para(doc, "Title: ______________________          Title: ______________________")
    doc_para(doc, "Signature: __________________          Signature: __________________")
    doc_para(doc, "Date: ______________________          Date: ______________________")

    out = TEMPLATES / "Synvoric-Master-Service-Agreement.docx"
    doc.save(str(out))
    print(f"Created {out}")


def build_contractor_agreement():
    doc = Document()
    doc_disclaimer(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("CONTRACTOR AGREEMENT")
    r.font.size = DocPt(18)
    r.font.bold = True
    r.font.color.rgb = DocRGB(0x05, 0x1C, 0x2C)

    doc_para(doc, "Effective Date: [EFFECTIVE DATE]")
    doc_para(doc, 'This Contractor Agreement ("Agreement") is between:')
    doc_para(doc, 'Synvoric ("Company")  ·  anshul.agrawal@synvoric.com  ·  www.synvoric.com', bold=True)
    doc_para(doc, 'Contractor: [CONTRACTOR FULL NAME] ("Contractor")', bold=True)
    doc_para(doc, "Address: [CONTRACTOR ADDRESS]  ·  PAN/Tax ID: [___________]")

    sections = [
        ("1. Engagement", [
            'Synvoric engages Contractor as an independent contractor to perform [DESCRIBE SERVICES — e.g., software development, frontend engineering] ("Services") for Synvoric clients as assigned by Synvoric.',
            "Contractor is not an employee. Contractor is responsible for all taxes, insurance, and statutory obligations.",
            "Contractor shall devote [FULL-TIME / minimum ___ hours per week] to assigned engagements unless otherwise agreed.",
        ]),
        ("2. Term", [
            "This Agreement begins on the Effective Date and continues until terminated under Section 11.",
            "Individual client assignments may be documented in assignment letters or SOWs referencing this Agreement.",
        ]),
        ("3. Payment Terms", [
            "Compensation: [INR/USD] [AMOUNT] per [month / hour / day] for Services performed.",
            "Contractor shall submit timesheets or invoices by [DATE EACH MONTH]. Payment shall be made within [15/30] days of Synvoric's receipt of a valid invoice.",
            "Synvoric may withhold payment for disputed hours pending resolution, with written explanation.",
            "No benefits, leave, or employment entitlements are provided unless required by applicable law.",
            "Contractor is responsible for all income tax, GST, and professional tax obligations arising from payments.",
        ]),
        ("4. Confidentiality & Non-Disclosure", [
            '"Confidential Information" includes Synvoric and client business information, source code, data, pricing, customer lists, and any information marked or reasonably understood as confidential.',
            "Contractor shall: (a) not disclose Confidential Information to any third party; (b) use it only to perform Services; (c) protect it with at least the same care used for Contractor's own confidential information.",
            "Contractor shall not discuss client engagements on social media or public forums without written approval.",
            "Upon termination, Contractor shall return or destroy all Confidential Information and certify destruction if requested.",
            "Confidentiality obligations survive termination indefinitely for trade secrets and for [3] years for other Confidential Information.",
            "Breach of this section may result in immediate termination and legal remedies including injunctive relief.",
        ]),
        ("5. Intellectual Property Ownership", [
            'All work product, code, designs, documentation, inventions, and materials created by Contractor in connection with Services ("Work Product") shall be the exclusive property of Synvoric (or Synvoric\'s client as directed by Synvoric).',
            "Work Product shall be deemed 'work made for hire' to the extent permitted by law. To the extent it is not, Contractor hereby irrevocably assigns all right, title, and interest in Work Product to Synvoric.",
            "Contractor retains no license to use Work Product except as authorized during the engagement.",
            "Contractor shall execute further documents reasonably required to perfect Synvoric's IP rights.",
            "Contractor represents that Work Product will not infringe third-party IP and has disclosed any pre-existing IP that may be incorporated (Pre-Existing IP).",
            "Pre-Existing IP incorporated into Work Product grants Synvoric a perpetual, royalty-free, worldwide license to use, modify, and sublicense such Pre-Existing IP as part of the Work Product.",
        ]),
        ("6. Non-Solicitation", [
            "During the term and for [12] months after termination, Contractor shall not, directly or indirectly:",
            "(a) solicit or encourage any Synvoric client to terminate or reduce business with Synvoric;",
            "(b) provide services directly to any Synvoric client that Contractor served during the [12] months prior to termination, without Synvoric's written consent;",
            "(c) solicit for employment or engagement any Synvoric employee or contractor.",
            "General advertising not targeted at specific Synvoric clients or personnel is permitted.",
            "Breach of this section entitles Synvoric to injunctive relief and liquidated damages of [INR/USD AMOUNT] per violation, without prejudice to other remedies.",
        ]),
        ("7. Conduct & Compliance", [
            "Contractor shall comply with Synvoric policies, client policies, and applicable laws including data protection requirements.",
            "Contractor shall not accept gifts or payments from clients without Synvoric's disclosure and approval.",
            "Contractor warrants no conflict of interest prevents full performance of Services.",
        ]),
        ("8. Representations", [
            "Contractor has the skills, qualifications, and legal right to perform Services.",
            "Contractor is not bound by any agreement that would prevent performance of this Agreement.",
            "Contractor shall maintain professional indemnity or errors & omissions insurance if required for the engagement.",
        ]),
        ("9. Limitation of Liability", [
            "Contractor's liability to Synvoric shall not exceed fees received in the [3] months preceding the claim, except for breaches of confidentiality, IP, or non-solicitation.",
        ]),
        ("10. Termination", [
            "Either Party may terminate with [15/30] days written notice.",
            "Synvoric may terminate immediately for breach of confidentiality, IP, non-solicitation, misconduct, or failure to perform.",
            "Upon termination, Contractor shall deliver all Work Product in progress and transition knowledge to Synvoric or its designee.",
            "Sections 4–6 and payment obligations for Services rendered survive termination.",
        ]),
        ("11. Governing Law", [
            "This Agreement is governed by the laws of [India / State of ____________].",
            "Disputes shall be subject to the exclusive jurisdiction of courts in [City, State/Country].",
        ]),
    ]

    for heading, paragraphs in sections:
        doc_heading(doc, heading, 2)
        for para in paragraphs:
            if para.startswith("(a)") or para.startswith("(b)") or para.startswith("(c)"):
                doc_bullet(doc, para)
            else:
                doc_para(doc, para)

    doc.add_paragraph()
    doc_heading(doc, "Signatures", 2)
    doc_para(doc, "SYNVORIC                              CONTRACTOR")
    doc_para(doc, "Name: _______________________          Name: _______________________")
    doc_para(doc, "Title: ______________________          Signature: __________________")
    doc_para(doc, "Signature: __________________          Date: ______________________")
    doc_para(doc, "Date: ______________________")

    out = TEMPLATES / "Synvoric-Contractor-Agreement.docx"
    doc.save(str(out))
    print(f"Created {out}")


def build_readme():
    readme = ASSETS / "README.md"
    readme.write_text("""# Synvoric Assets

Business documents and capability materials for Synvoric.

## Capability Deck

| File | Description |
|------|-------------|
| `Synvoric-Capability-Deck.pptx` | 12-slide McKinsey-style capability presentation |

**Also available online:** https://synvoric.com/deck (browser view + Print to PDF)

## Legal & Commercial Templates

Located in `templates/`. **Review with legal counsel before use.**

| File | Purpose |
|------|---------|
| `Synvoric-Client-Proposal-Template.docx` | Client-facing proposal for staffing engagements |
| `Synvoric-Master-Service-Agreement.docx` | MSA covering confidentiality, IP, payment, non-solicitation |
| `Synvoric-Contractor-Agreement.docx` | Independent contractor agreement with NDA, IP assignment, payment terms, non-solicitation |

## Regenerating Files

```bash
python3 scripts/generate_assets.py
```

## Placeholders

All templates use `[BRACKET]` placeholders — replace before sending to clients or contractors.
""")
    print(f"Created {readme}")


if __name__ == "__main__":
    ASSETS.mkdir(parents=True, exist_ok=True)
    TEMPLATES.mkdir(parents=True, exist_ok=True)
    build_ppt()
    build_proposal()
    build_msa()
    build_contractor_agreement()
    build_readme()
    print("Done.")
