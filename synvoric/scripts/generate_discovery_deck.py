#!/usr/bin/env python3
"""
Synvoric × Consumables Solutions — Discovery Deck (first-call version).

Purpose: capability-led first discussion. No named candidates.
Indicative pricing as ranges.
"""

from __future__ import annotations

from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
LOGO = ROOT / "public" / "logo.png"
OUT = ASSETS / "Synvoric-Consumables-Discovery-Deck.pptx"

NAVY = RGBColor(0x05, 0x1C, 0x2C)
BLUE = RGBColor(0x00, 0x66, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
BODY = RGBColor(0x33, 0x41, 0x55)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
DARK_NAVY = RGBColor(0x03, 0x14, 0x22)
NAVY_LIGHT = RGBColor(0x1A, 0x2E, 0x45)

ML = Inches(0.72)
CW = Inches(8.56)
BODY_Y = Inches(1.70)

CLIENT = "Consumables Solutions"
TOTAL_SLIDES = 17
SOURCE_LINE = "Source: Synvoric internal data, FY 2025–26"


# ── helpers ────────────────────────────────────────────────────────────

def blank_layout(prs):
    return prs.slide_layouts[6]


def set_bg(slide, rgb: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def txt(slide, left, top, width, height, text, size=11, bold=False, color=BODY,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tf


def bullets(tf, items, size=10.5, color=BODY, spacing=6):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.bullet = True


def style_cell(cell, fill=None, margin=0.10):
    cell.margin_left = Inches(margin)
    cell.margin_right = Inches(margin)
    cell.margin_top = Inches(margin * 0.5)
    cell.margin_bottom = Inches(margin * 0.5)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def line_h(slide, x, y, w, color=BLUE, thickness=0.025):
    return rect(slide, x, y, w, Inches(thickness), color)


def footer(slide, n: int, dark: bool = False, source: str | None = None):
    if source and not dark:
        txt(slide, ML, Inches(5.05), CW, Inches(0.18), source, 8, color=SLATE, italic=True)
    rect(slide, Inches(0), Inches(5.28), Inches(10), Inches(0.28),
         DARK_NAVY if dark else NAVY)
    lc = RGBColor(0x94, 0xA3, 0xB8) if dark else SLATE
    bc = RGBColor(0x66, 0x99, 0xFF) if dark else BLUE
    txt(slide, ML, Inches(5.31), Inches(6), Inches(0.22),
        f"Synvoric  ·  Discovery call with {CLIENT}  ·  Confidential", 7.5, color=lc)
    txt(slide, Inches(7.5), Inches(5.31), Inches(1.8), Inches(0.22),
        f"{n:02d} / {TOTAL_SLIDES:02d}", 7.5, bold=True, color=bc, align=PP_ALIGN.RIGHT)


def header(slide, eyebrow: str, action_title: str):
    txt(slide, ML, Inches(0.42), CW, Inches(0.22), eyebrow.upper(), 8.5, bold=True, color=BLUE)
    txt(slide, ML, Inches(0.68), CW, Inches(0.85), action_title, 20, bold=True, color=NAVY)
    line_h(slide, ML, Inches(1.55), Inches(0.42), color=BLUE)


def divider(prs, num: str, title: str, subtitle: str, n: int):
    s = prs.slides.add_slide(blank_layout(prs))
    set_bg(s, NAVY)
    txt(s, ML, Inches(2.40), Inches(3), Inches(0.7), num, 52, bold=True, color=NAVY_LIGHT)
    txt(s, ML, Inches(3.00), CW, Inches(0.65), title, 30, bold=True, color=WHITE)
    line_h(s, ML, Inches(3.72), Inches(0.55), color=BLUE, thickness=0.03)
    txt(s, ML, Inches(3.85), CW, Inches(0.35), subtitle, 12, color=RGBColor(0xCB, 0xD5, 0xE1))
    footer(s, n, dark=True)
    return s


# ── slides ─────────────────────────────────────────────────────────────

def cover(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    set_bg(s, NAVY)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(3.05), Inches(0.80), width=Inches(3.9))
    txt(s, Inches(1), Inches(2.50), Inches(8), Inches(0.3),
        "AI · DATA · ENGINEERING · SOLUTIONS", 10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.00), Inches(8), Inches(0.55),
        "Capability & Engagement Overview", 24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.65), Inches(8), Inches(0.35),
        f"Initial discussion with {CLIENT}", 14,
        color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(4.15), Inches(8), Inches(0.25),
        f"{date.today().strftime('%B %Y')}  ·  synvoric.com", 9, color=SLATE, align=PP_ALIGN.CENTER)
    footer(s, 1, dark=True)


def executive_summary(prs, n=2):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Executive Summary",
           "Synvoric can deliver design and frontend capacity for Consumables Solutions")

    blocks = [
        ("Capability",
         "Synvoric maintains a vetted bench of UX/UI designers and frontend developers across "
         "junior, mid, and senior experience tiers — available on flexible engagement terms."),
        ("Approach",
         "We propose to share matched candidate profiles aligned to your specific role briefs, "
         "stack, and seniority needs — and let you interview before commitment."),
        ("Commercials",
         "Indicative monthly rates start from $3,000 USD per resource and scale by role and "
         "seniority. Final rate confirmed once role briefs are shared."),
    ]
    for i, (head, body) in enumerate(blocks):
        y = BODY_Y + Inches(i * 1.05)
        rect(s, ML, y, Inches(0.05), Inches(0.85), BLUE)
        txt(s, ML + Inches(0.18), y, Inches(1.6), Inches(0.30), head.upper(), 10, bold=True, color=BLUE)
        txt(s, ML + Inches(1.85), y, Inches(6.50), Inches(0.95), body, 11, color=BODY)
    footer(s, n, source=SOURCE_LINE)


def agenda(prs, n=3):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Agenda", "What we will cover today")
    items = [
        ("About Synvoric", "Who we are and why staffing leaders trust us"),
        ("Our Capabilities", "Talent areas, tech stack, and bench depth"),
        ("Engagement Models", "Staff augmentation, dedicated teams, project delivery"),
        ("Indicative Commercials", "Rate ranges and how pricing is determined"),
        ("Next Steps", "Path to candidate profiles and engagement"),
    ]
    for i, (head, sub) in enumerate(items):
        y = BODY_Y + Inches(i * 0.62)
        txt(s, ML, y, Inches(0.45), Inches(0.40), f"{i + 1:02d}", 14, bold=True, color=BLUE)
        txt(s, Inches(1.25), y, Inches(3.0), Inches(0.40), head, 12.5, bold=True, color=NAVY)
        txt(s, Inches(4.30), y, Inches(4.30), Inches(0.40), sub, 11, color=SLATE)
    footer(s, n)


def about_synvoric(prs, n=4):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "About Synvoric",
           "A focused staffing partner for design, frontend, and AI engineering talent")
    txt(s, ML, BODY_Y, Inches(4.10), Inches(1.05),
        "Synvoric is a professional IT staffing firm specializing in frontend developers, "
        "UX/UI designers, and AI engineers. We help companies scale product teams with "
        "pre-vetted talent — integrated into your tools, sprints, and delivery cadence.",
        11, color=BODY)

    stats = [("50+", "Developers Placed"), ("10+", "Active Clients"),
             ("100%", "Client Retention"), ("2 Wks", "Avg. Time to Hire")]
    for i, (val, lbl) in enumerate(stats):
        x = ML + Inches(i * 2.15)
        box = rect(s, x, Inches(3.10), Inches(1.95), Inches(1.10), LIGHT_BG, line=BORDER)
        btf = box.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = btf.paragraphs[0]
        vp.text = val
        vp.font.size = Pt(26)
        vp.font.bold = True
        vp.font.color.rgb = BLUE
        vp.alignment = PP_ALIGN.CENTER
        lp = btf.add_paragraph()
        lp.text = lbl
        lp.font.size = Pt(8)
        lp.font.color.rgb = SLATE
        lp.alignment = PP_ALIGN.CENTER
    footer(s, n, source=SOURCE_LINE)


def why_synvoric(prs, n=5):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Differentiators",
           "Four reasons leading product teams choose Synvoric")
    pillars = [
        ("Rigorous Vetting", "Portfolio review, technical assessments, and live interviews.",
         "100% client retention"),
        ("Speed to Deploy", "Candidates onboarded inside two weeks of confirmation.",
         "2 weeks average"),
        ("Flexible Engagement", "Scale up or down with 30-day notice — no long-term lock-in.",
         "30-day rolling notice"),
        ("Replacement Guarantee", "Replace any resource within 30 days at no extra cost.",
         "Zero risk to start"),
    ]
    for i, (title, body, proof) in enumerate(pillars):
        col, row = i % 2, i // 2
        x = ML + Inches(col * 4.30)
        y = BODY_Y + Inches(row * 1.65)
        rect(s, x, y, Inches(0.05), Inches(1.45), BLUE)
        txt(s, x + Inches(0.18), y, Inches(4.0), Inches(0.30), title, 12.5, bold=True, color=NAVY)
        txt(s, x + Inches(0.18), y + Inches(0.32), Inches(4.0), Inches(0.65), body, 10.5, color=BODY)
        txt(s, x + Inches(0.18), y + Inches(1.05), Inches(4.0), Inches(0.30),
            proof, 9.5, bold=True, color=BLUE, italic=True)
    footer(s, n, source=SOURCE_LINE)


def capabilities(prs, n=7):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Capabilities",
           "What we staff — frontend, design, and full-stack product talent")
    caps = [
        ("Frontend Developers",
         "React, Next.js, Angular, Vue, TypeScript — junior to senior tiers."),
        ("UX / UI Designers",
         "Figma-led design — user research, wireframes, prototypes, design systems."),
        ("Full-Stack & Mobile",
         "MERN / MEAN, React Native, Flutter — full product capability when needed."),
        ("AI / ML Engineers",
         "LLM integration, ML pipelines, MLOps — for AI-powered product features."),
    ]
    for i, (title, body) in enumerate(caps):
        col, row = i % 2, i // 2
        x = ML + Inches(col * 4.30)
        y = BODY_Y + Inches(row * 1.65)
        box = rect(s, x, y, Inches(4.10), Inches(1.45), LIGHT_BG, line=BORDER)
        line_h(s, x, y, Inches(4.10), color=BLUE, thickness=0.04)
        txt(s, x + Inches(0.18), y + Inches(0.18), Inches(3.75), Inches(0.30),
            title, 12.5, bold=True, color=NAVY)
        txt(s, x + Inches(0.18), y + Inches(0.55), Inches(3.75), Inches(0.85),
            body, 10.5, color=BODY)
    footer(s, n)


def tech_stack(prs, n=8):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Technology Stack",
           "Technologies our developers and designers master")
    groups = [
        ("Frontend",
         "React · Next.js · Angular · Vue.js · TypeScript · Tailwind · Redux · Vite"),
        ("Design",
         "Figma · Adobe XD · Sketch · Framer · Design Systems · Prototyping"),
        ("Mobile",
         "React Native · Flutter · Swift · Kotlin · iOS · Android · Expo · Firebase"),
        ("Backend & Cloud",
         "Node.js · MongoDB · PostgreSQL · GraphQL · AWS · Docker · Kubernetes"),
    ]
    for i, (cat, items) in enumerate(groups):
        y = BODY_Y + Inches(i * 0.78)
        rect(s, ML, y, Inches(0.05), Inches(0.62), BLUE)
        txt(s, ML + Inches(0.18), y, Inches(2.0), Inches(0.30),
            cat.upper(), 10, bold=True, color=BLUE)
        txt(s, ML + Inches(0.18), y + Inches(0.30), Inches(8.0), Inches(0.32),
            items, 11, color=BODY)
    footer(s, n)


def talent_pool(prs, n=9):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Talent Pool",
           "Bench depth across roles and seniority — matched to your specific brief")

    # Top: depth summary cards
    depths = [
        ("Frontend Developers", "25+", "Junior · Mid · Senior"),
        ("UX / UI Designers", "10+", "Junior · Mid · Senior"),
        ("Full-Stack Engineers", "15+", "MERN · MEAN · Cloud"),
        ("AI / ML Engineers", "8+", "ML · LLMs · MLOps"),
    ]
    for i, (role, count, tiers) in enumerate(depths):
        x = ML + Inches(i * 2.15)
        box = rect(s, x, BODY_Y, Inches(1.95), Inches(1.20), LIGHT_BG, line=BORDER)
        line_h(s, x, BODY_Y, Inches(1.95), color=BLUE, thickness=0.04)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.12), Inches(1.70), Inches(0.30),
            count, 18, bold=True, color=BLUE)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.45), Inches(1.70), Inches(0.30),
            role, 10, bold=True, color=NAVY)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.78), Inches(1.70), Inches(0.30),
            tiers, 8.5, color=SLATE)

    # Bottom: seniority tiers explanation
    tier_y = Inches(3.10)
    txt(s, ML, tier_y, CW, Inches(0.25), "SENIORITY TIERS", 9, bold=True, color=BLUE)
    tiers_data = [
        ("Junior", "1–2 years", "Mentored execution on well-defined tasks"),
        ("Mid-level", "3–5 years", "Independent delivery, code reviews, design ownership"),
        ("Senior", "5+ years", "Architecture, mentoring, technical leadership"),
    ]
    for i, (tier, exp, desc) in enumerate(tiers_data):
        x = ML + Inches(i * 2.86)
        y = tier_y + Inches(0.30)
        box = rect(s, x, y, Inches(2.70), Inches(1.30), WHITE, line=BORDER)
        txt(s, x + Inches(0.15), y + Inches(0.12), Inches(2.40), Inches(0.30),
            tier, 11, bold=True, color=NAVY)
        txt(s, x + Inches(0.15), y + Inches(0.42), Inches(2.40), Inches(0.25),
            exp, 9.5, bold=True, color=BLUE)
        txt(s, x + Inches(0.15), y + Inches(0.72), Inches(2.40), Inches(0.55),
            desc, 9.5, color=BODY)
    footer(s, n, source=SOURCE_LINE)


def engagement_models(prs, n=11):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Engagement Models",
           "Three ways to engage — chosen based on your scope and team needs")
    models = [
        ("Staff Augmentation",
         "Individual developers or designers embedded in your team.",
         "Best for: filling specific skill gaps quickly"),
        ("Dedicated Team",
         "A pod of 2–5 resources working exclusively on your product.",
         "Best for: scaling a product roadmap"),
        ("Project Delivery",
         "End-to-end ownership — discovery through deployment.",
         "Best for: defined-scope projects with deadlines"),
    ]
    for i, (title, body, fit) in enumerate(models):
        x = ML + Inches(i * 2.86)
        box = rect(s, x, BODY_Y, Inches(2.70), Inches(2.85), LIGHT_BG, line=BORDER)
        line_h(s, x, BODY_Y, Inches(2.70), color=BLUE, thickness=0.04)
        txt(s, x + Inches(0.15), BODY_Y + Inches(0.18), Inches(2.40), Inches(0.32),
            title, 12.5, bold=True, color=NAVY)
        txt(s, x + Inches(0.15), BODY_Y + Inches(0.62), Inches(2.40), Inches(1.10),
            body, 10.5, color=BODY)
        # Best for tag
        tag_y = BODY_Y + Inches(1.95)
        line_h(s, x + Inches(0.15), tag_y, Inches(2.40), color=BORDER, thickness=0.01)
        txt(s, x + Inches(0.15), tag_y + Inches(0.10), Inches(2.40), Inches(0.22),
            "BEST FOR", 8, bold=True, color=BLUE)
        txt(s, x + Inches(0.15), tag_y + Inches(0.35), Inches(2.40), Inches(0.50),
            fit.replace("Best for: ", ""), 10, color=BODY)
    footer(s, n)


def engagement_structure(prs, n=12):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Engagement Structure",
           "Clear ownership — one account manager, dedicated specialists")

    cx = Inches(5.0)

    client_box = rect(s, cx - Inches(1.6), Inches(1.85), Inches(3.2), Inches(0.70),
                      NAVY, line=BORDER)
    cb = client_box.text_frame
    cb.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = cb.paragraphs[0]
    cp.text = f"{CLIENT} — Project Owner"
    cp.font.size = Pt(11.5)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER

    am_box = rect(s, cx - Inches(1.6), Inches(3.00), Inches(3.2), Inches(0.70),
                  BLUE, line=BORDER)
    ab = am_box.text_frame
    ab.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = ab.paragraphs[0]
    ap.text = "Synvoric — Account Manager"
    ap.font.size = Pt(11.5)
    ap.font.bold = True
    ap.font.color.rgb = WHITE
    ap.alignment = PP_ALIGN.CENTER

    res_y = Inches(4.20)
    for i, role in enumerate(["UX / UI Designer", "Frontend Developer"]):
        x = cx + Inches(-1.85 if i == 0 else 0.05)
        box = rect(s, x, res_y, Inches(1.80), Inches(0.70), LIGHT_BG, line=BORDER)
        bt = box.text_frame
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bt.margin_left = Inches(0.10)
        np = bt.paragraphs[0]
        np.text = "Dedicated Resource"
        np.font.size = Pt(9)
        np.font.color.rgb = SLATE
        np.alignment = PP_ALIGN.CENTER
        rp = bt.add_paragraph()
        rp.text = role
        rp.font.size = Pt(10.5)
        rp.font.bold = True
        rp.font.color.rgb = NAVY
        rp.alignment = PP_ALIGN.CENTER

    # connectors
    rect(s, cx - Inches(0.01), Inches(2.55), Inches(0.02), Inches(0.45), SLATE)
    rect(s, cx - Inches(0.95), Inches(3.90), Inches(1.90), Inches(0.02), SLATE)
    rect(s, cx - Inches(0.01), Inches(3.70), Inches(0.02), Inches(0.20), SLATE)
    rect(s, cx - Inches(0.95), Inches(3.92), Inches(0.02), Inches(0.28), SLATE)
    rect(s, cx + Inches(0.95), Inches(3.92), Inches(0.02), Inches(0.28), SLATE)

    notes = [
        ("Weekly cadence", "Status report every Friday"),
        ("Single point of contact", "All escalations via Account Manager"),
        ("No team management overhead", "Synvoric handles payroll and oversight"),
    ]
    for i, (h, b) in enumerate(notes):
        y = Inches(1.85 + i * 1.05)
        txt(s, ML, y, Inches(2.0), Inches(0.25), h, 10, bold=True, color=NAVY)
        txt(s, ML, y + Inches(0.27), Inches(2.0), Inches(0.55), b, 9.5, color=SLATE)
    footer(s, n)


def indicative_pricing(prs, n=14):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Indicative Commercials",
           "Rate ranges by role — final pricing confirmed once we see your role brief")

    rows = [
        ["Role", "Experience", "Allocation", "Indicative Range (USD / month)"],
        ["UX / UI Designer", "3 – 6 years", "Full-time", "$ 3,000  –  $ 4,500"],
        ["Frontend Developer", "3 – 6 years", "Full-time", "$ 3,500  –  $ 5,000"],
        ["Senior Specialist", "6+ years",     "Full-time", "On request"],
    ]
    row_heights = [Inches(0.42), Inches(0.50), Inches(0.50), Inches(0.50)]
    table_h = sum(row_heights, Inches(0))
    tbl = s.shapes.add_table(len(rows), 4, ML, BODY_Y, CW, table_h).table
    widths = [Inches(2.30), Inches(1.80), Inches(1.70), Inches(2.76)]
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = w
    for ri, h in enumerate(row_heights):
        tbl.rows[ri].height = h
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            style_cell(cell, fill=NAVY if ri == 0 else None)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(9.5 if ri == 0 else 10.5)
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else BODY
                if ci in (1, 2, 3):
                    p.alignment = PP_ALIGN.CENTER

    # Factors affecting rate
    factors_y = BODY_Y + table_h + Inches(0.25)
    txt(s, ML, factors_y, CW, Inches(0.25),
        "RATE IS DETERMINED BY", 9, bold=True, color=BLUE)
    factors = [
        ("Seniority", "Years of experience and depth of expertise"),
        ("Stack specialization", "Niche frameworks or tools command a premium"),
        ("Engagement length", "Longer commitments unlock preferred rates"),
        ("Allocation", "Full-time vs. part-time vs. shared resource"),
    ]
    for i, (k, v) in enumerate(factors):
        col, row = i % 2, i // 2
        x = ML + Inches(col * 4.30)
        y = factors_y + Inches(0.30 + row * 0.32)
        txt(s, x, y, Inches(1.30), Inches(0.28), k, 10, bold=True, color=NAVY)
        txt(s, x + Inches(1.30), y, Inches(2.85), Inches(0.28), v, 10, color=BODY)
    footer(s, n,
           source="Indicative ranges only. Final rate confirmed post role-brief alignment. All fees exclusive of applicable taxes.")


def process(prs, n=15):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "How We Work Together",
           "From discovery call to deployed team in under two weeks")
    steps = [
        ("Step 1", "Role Brief", "Share required roles, skills, seniority, and timeline"),
        ("Step 2", "Profiles", "Receive 2–3 matched candidate profiles within 3 days"),
        ("Step 3", "Interviews", "Interview shortlisted candidates at your convenience"),
        ("Step 4", "Onboarding", "Resources onboarded within 5–7 business days of sign-off"),
    ]
    for i, (phase, title, desc) in enumerate(steps):
        x = ML + Inches(i * 2.15)
        box = rect(s, x, BODY_Y, Inches(1.95), Inches(2.25), LIGHT_BG, line=BORDER)
        line_h(s, x, BODY_Y, Inches(1.95), color=BLUE, thickness=0.04)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.12), Inches(1.70), Inches(0.25),
            phase.upper(), 9, bold=True, color=BLUE)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.40), Inches(1.70), Inches(0.32),
            title, 12, bold=True, color=NAVY)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.80), Inches(1.70), Inches(1.30),
            desc, 10, color=BODY)
    footer(s, n)


def next_steps_closing(prs, n=16):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Next Steps",
           "What we need from Consumables Solutions to share matched profiles")
    asks = [
        ("Job Description",
         "Detailed JD for each role — responsibilities, skills, and any must-haves"),
        ("Indicative Start and End Date",
         "Target onboarding date and expected engagement duration"),
    ]
    txt(s, ML, BODY_Y, CW, Inches(0.30),
        "Once we have these inputs, expect matched profiles within 3 business days.",
        11, color=BODY)

    for i, (head, body) in enumerate(asks):
        y = BODY_Y + Inches(0.55 + i * 1.10)
        rect(s, ML, y, Inches(0.05), Inches(0.85), BLUE)
        txt(s, ML + Inches(0.18), y, Inches(2.6), Inches(0.32),
            head.upper(), 11, bold=True, color=NAVY)
        txt(s, ML + Inches(0.18), y + Inches(0.36), Inches(7.8), Inches(0.55),
            body, 10.5, color=BODY)

    cta = rect(s, ML, Inches(4.55), CW, Inches(0.45), NAVY)
    ctf = cta.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.20)
    cp = ctf.paragraphs[0]
    cp.text = "  Contact: anshul.agrawal@synvoric.com  ·  www.synvoric.com"
    cp.font.size = Pt(11)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    footer(s, n)


def thank_you_qa(prs, n=17):
    s = prs.slides.add_slide(blank_layout(prs))
    set_bg(s, NAVY)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(3.55), Inches(0.85), width=Inches(2.9))
    txt(s, Inches(1), Inches(2.45), Inches(8), Inches(0.5),
        "Thank You", 36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    line_h(s, Inches(4.65), Inches(3.10), Inches(0.70), color=BLUE, thickness=0.04)
    txt(s, Inches(1), Inches(3.25), Inches(8), Inches(0.40),
        "Questions & Discussion", 18, bold=True,
        color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.85), Inches(8), Inches(0.30),
        "We look forward to partnering with Consumables Solutions", 11,
        color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(4.30), Inches(8), Inches(0.30),
        "anshul.agrawal@synvoric.com  ·  synvoric.com", 11,
        bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    footer(s, n, dark=True)


# ── build ──────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    cover(prs)
    executive_summary(prs, 2)
    agenda(prs, 3)
    about_synvoric(prs, 4)
    why_synvoric(prs, 5)
    divider(prs, "01", "Our Capability",
            "Talent areas, technology stack, and bench depth", 6)
    capabilities(prs, 7)
    tech_stack(prs, 8)
    talent_pool(prs, 9)
    divider(prs, "02", "The Engagement",
            "How we engage and how the work flows", 10)
    engagement_models(prs, 11)
    engagement_structure(prs, 12)
    divider(prs, "03", "Commercials & Next Steps",
            "Indicative pricing and what we need from you", 13)
    indicative_pricing(prs, 14)
    process(prs, 15)
    next_steps_closing(prs, 16)
    thank_you_qa(prs, 17)

    prs.save(str(OUT))
    print(f"Created {OUT} ({TOTAL_SLIDES} slides)")


if __name__ == "__main__":
    ASSETS.mkdir(parents=True, exist_ok=True)
    build()
