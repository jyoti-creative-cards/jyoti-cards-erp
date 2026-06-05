#!/usr/bin/env python3
"""Synvoric × Consumables Solutions — Proposal Deck v2 (McKinsey style)."""

from __future__ import annotations

from pathlib import Path
from datetime import date

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"
LOGO = ROOT / "public" / "logo.png"
OUT = ASSETS / "Synvoric-Consumables-Solutions-Proposal-Deck.pptx"
PHOTOS = {
    "Ayush Shukla": ASSETS / "photos" / "ayush.jpg",
    "Sarthak Shukla": ASSETS / "photos" / "sarthak.jpg",
}

# Palette
NAVY = RGBColor(0x05, 0x1C, 0x2C)
BLUE = RGBColor(0x00, 0x66, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x64, 0x74, 0x8B)
BODY = RGBColor(0x33, 0x41, 0x55)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
DARK_NAVY = RGBColor(0x03, 0x14, 0x22)
NAVY_LIGHT = RGBColor(0x1A, 0x2E, 0x45)

# Layout
ML = Inches(0.72)
CW = Inches(8.56)
BODY_Y = Inches(1.70)

CLIENT = "Consumables Solutions"
RATE_USD = 3500
ROLES = 2
MONTHLY_TOTAL = RATE_USD * ROLES
TOTAL_SLIDES = 18
SOURCE_LINE = "Source: Synvoric internal data, FY 2025–26"


# ── primitives ────────────────────────────────────────────────────────

def blank_layout(prs):
    return prs.slide_layouts[6]


def set_bg(slide, rgb: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def txt(slide, left, top, width, height, text, size=11, bold=False, color=BODY, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tf


def bullets(tf, items, size=10.5, color=BODY, spacing=6):
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.bullet = True


def style_cell(cell, fill=None, margin=0.10, vertical=MSO_ANCHOR.MIDDLE):
    cell.margin_left = Inches(margin)
    cell.margin_right = Inches(margin)
    cell.margin_top = Inches(margin * 0.5)
    cell.margin_bottom = Inches(margin * 0.5)
    cell.vertical_anchor = vertical
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def line_h(slide, x, y, w, color=BLUE, thickness=0.025):
    return rect(slide, x, y, w, Inches(thickness), color)


def footer(slide, n: int, dark=False, source: str | None = None):
    # Source line above footer bar
    if source and not dark:
        txt(slide, ML, Inches(5.05), CW, Inches(0.18), source, 8, color=SLATE, italic=True)

    bar = rect(slide, Inches(0), Inches(5.28), Inches(10), Inches(0.28),
               DARK_NAVY if dark else NAVY)
    lc = RGBColor(0x94, 0xA3, 0xB8) if dark else SLATE
    bc = RGBColor(0x66, 0x99, 0xFF) if dark else BLUE
    txt(slide, ML, Inches(5.31), Inches(6), Inches(0.22),
        f"Synvoric  ·  Prepared for {CLIENT}  ·  Confidential", 7.5, color=lc)
    txt(slide, Inches(7.5), Inches(5.31), Inches(1.8), Inches(0.22),
        f"{n:02d} / {TOTAL_SLIDES:02d}", 7.5, bold=True, color=bc, align=PP_ALIGN.RIGHT)


def header(slide, eyebrow: str, action_title: str):
    txt(slide, ML, Inches(0.42), CW, Inches(0.22), eyebrow.upper(), 8.5, bold=True, color=BLUE)
    txt(slide, ML, Inches(0.68), CW, Inches(0.85), action_title, 20, bold=True, color=NAVY)
    line_h(slide, ML, Inches(1.55), Inches(0.42), color=BLUE)


def divider(prs, section_num: str, title: str, subtitle: str, n: int):
    s = prs.slides.add_slide(blank_layout(prs))
    set_bg(s, NAVY)
    txt(s, ML, Inches(2.40), Inches(3), Inches(0.7), section_num, 52, bold=True, color=NAVY_LIGHT)
    txt(s, ML, Inches(3.00), CW, Inches(0.65), title, 30, bold=True, color=WHITE)
    line_h(s, ML, Inches(3.72), Inches(0.55), color=BLUE, thickness=0.03)
    txt(s, ML, Inches(3.85), CW, Inches(0.35), subtitle, 12, color=RGBColor(0xCB, 0xD5, 0xE1))
    footer(s, n, dark=True)
    return s


# ── slides ────────────────────────────────────────────────────────────

def cover(prs):
    s = prs.slides.add_slide(blank_layout(prs))
    set_bg(s, NAVY)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(3.05), Inches(0.80), width=Inches(3.9))
    txt(s, Inches(1), Inches(2.50), Inches(8), Inches(0.3),
        "AI · DATA · ENGINEERING · SOLUTIONS", 10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.00), Inches(8), Inches(0.55),
        "Staffing Proposal", 26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.60), Inches(8), Inches(0.35),
        f"Prepared exclusively for {CLIENT}", 14, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(4.10), Inches(8), Inches(0.25),
        f"{date.today().strftime('%B %Y')}  ·  synvoric.com", 9, color=SLATE, align=PP_ALIGN.CENTER)
    footer(s, 1, dark=True)


def executive_summary(prs, n=2):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Executive Summary",
           "Dedicated 2-person design + frontend team, deployed in 2 weeks at $7,000 / month")

    blocks = [
        ("Situation",
         f"{CLIENT} needs additional design and frontend engineering capacity to accelerate "
         "digital product delivery without expanding fixed headcount."),
        ("Solution",
         "Synvoric proposes a dedicated 2-person team — a Senior UX/UI Designer (Ayush Shukla) "
         "and a Frontend Developer (Sarthak Shukla) — embedded into your workflow."),
        ("Commercials",
         f"$3,500 USD / month per resource · $7,000 / month total · 30-day rolling notice · "
         "replacement guarantee · onboarded within 5–7 business days of signing."),
    ]
    for i, (head, body) in enumerate(blocks):
        y = BODY_Y + Inches(i * 1.05)
        bar = rect(s, ML, y, Inches(0.05), Inches(0.85), BLUE)
        txt(s, ML + Inches(0.18), y, Inches(1.6), Inches(0.30), head.upper(), 10, bold=True, color=BLUE)
        txt(s, ML + Inches(1.85), y, Inches(6.50), Inches(0.95), body, 11, color=BODY)
    footer(s, n, source=SOURCE_LINE)


def agenda(prs, n=3):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Overview", "What we will cover today")
    items = [
        ("About Synvoric", "Who we are and why staffing leaders trust us"),
        ("The Proposal", "Understanding your need and our recommended team"),
        ("Candidate Profiles", "Meet Ayush (Designer) and Sarthak (Frontend Dev)"),
        ("Commercials & Risk", "Pricing, terms, and built-in safeguards"),
        ("Next Steps", "Path to a deployed team within one week"),
    ]
    for i, (head, sub) in enumerate(items):
        y = BODY_Y + Inches(i * 0.62)
        txt(s, ML, y, Inches(0.45), Inches(0.40), f"{i + 1:02d}", 14, bold=True, color=BLUE)
        txt(s, Inches(1.25), y, Inches(3.0), Inches(0.40), head, 12.5, bold=True, color=NAVY)
        txt(s, Inches(4.30), y, Inches(4.30), Inches(0.40), sub, 11, color=SLATE)
    footer(s, n)


def about_synvoric(prs, n=4):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "About Synvoric",
           "A focused staffing partner for design, frontend and AI engineering talent")
    txt(s, ML, BODY_Y, Inches(4.10), Inches(1.05),
        "Synvoric is a professional IT staffing firm specializing in frontend developers, "
        "UX/UI designers, and AI engineers. We help companies scale product teams with "
        "pre-vetted talent — integrated into your tools, sprints, and delivery cadence.",
        11, color=BODY)

    stats = [("50+", "Developers Placed"), ("10+", "Active Clients"),
             ("100%", "Client Retention"), ("2 Wks", "Avg. Time to Hire")]
    for i, (val, lbl) in enumerate(stats):
        x = ML + Inches(i * 2.15)
        box = rect(s, x, Inches(3.10), Inches(1.95), Inches(1.10), LIGHT_BG, line=BORDER)
        btf = box.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        vp = btf.paragraphs[0]
        vp.text = val
        vp.font.size = Pt(26)
        vp.font.bold = True
        vp.font.color.rgb = BLUE
        vp.alignment = PP_ALIGN.CENTER
        lp = btf.add_paragraph()
        lp.text = lbl
        lp.font.size = Pt(8)
        lp.font.color.rgb = SLATE
        lp.alignment = PP_ALIGN.CENTER
    footer(s, n, source=SOURCE_LINE)


def why_synvoric(prs, n=5):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Differentiators",
           "Four reasons leading product teams choose Synvoric")
    pillars = [
        ("Rigorous Vetting", "Portfolio review, technical assessments, and live interviews.",
         "100% client retention"),
        ("Speed to Deploy", "Candidates onboarded inside two weeks of confirmation.",
         "2 weeks average"),
        ("Flexible Engagement", "Scale up or down with 30-day notice — no long-term lock-in.",
         "30-day rolling notice"),
        ("Replacement Guarantee", "Replace any resource within 30 days at no extra cost.",
         "Zero risk to start"),
    ]
    for i, (title, body, proof) in enumerate(pillars):
        col, row = i % 2, i // 2
        x = ML + Inches(col * 4.30)
        y = BODY_Y + Inches(row * 1.65)
        bar = rect(s, x, y, Inches(0.05), Inches(1.45), BLUE)
        txt(s, x + Inches(0.18), y, Inches(4.0), Inches(0.30), title, 12.5, bold=True, color=NAVY)
        txt(s, x + Inches(0.18), y + Inches(0.32), Inches(4.0), Inches(0.65), body, 10.5, color=BODY)
        txt(s, x + Inches(0.18), y + Inches(1.05), Inches(4.0), Inches(0.30),
            proof, 9.5, bold=True, color=BLUE, italic=True)
    footer(s, n, source=SOURCE_LINE)


def client_context(prs, n=7):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Client Context",
           f"{CLIENT} needs design + frontend capacity to accelerate digital product delivery")

    txt(s, ML, BODY_Y, CW, Inches(0.55),
        "Based on our discussions, we understand the following requirements:",
        11, color=BODY)
    tf = txt(s, ML, BODY_Y + Inches(0.55), CW, Inches(2.4), "", 11)
    bullets(tf, [
        "UX/UI Designer — Figma-led user flows, wireframes, and high-fidelity designs",
        "Frontend Developer — React / Next.js implementation from design to production",
        "Remote-ready resources with proven product delivery experience",
        "Immediate start with Synvoric-managed onboarding and oversight",
        "Flexible commercials with the ability to scale the team as projects evolve",
    ], size=11, spacing=8)
    footer(s, n)


def proposed_team(prs, n=8):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Proposed Solution",
           "A 2-person dedicated team purpose-built for your roadmap")

    rows = [
        ["#", "Role", "Resource", "Experience", "Key Skills"],
        ["1", "UX / UI Designer", "Ayush Shukla", "2+ yrs",
         "Figma · Wireframing · Prototyping · Design Systems"],
        ["2", "Frontend Developer", "Sarthak Shukla", "4 yrs",
         "React.js · Next.js · TypeScript · React Native"],
    ]
    row_heights = [Inches(0.42), Inches(0.55), Inches(0.55)]
    tbl = s.shapes.add_table(len(rows), 5, ML, BODY_Y,
                             CW, sum(row_heights, Inches(0))).table
    widths = [Inches(0.40), Inches(1.65), Inches(1.65), Inches(0.85), Inches(4.01)]
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = w
    for ri, h in enumerate(row_heights):
        tbl.rows[ri].height = h
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            style_cell(cell, fill=NAVY if ri == 0 else None)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(9.5 if ri == 0 else 10.5)
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else BODY
                if ci in (0, 3):
                    p.alignment = PP_ALIGN.CENTER

    txt(s, ML, Inches(3.60), CW, Inches(0.85),
        "Both resources work as dedicated Synvoric contractors under a single engagement. "
        "Includes weekly status reports, dedicated account manager, and replacement guarantee.",
        10.5, color=SLATE)
    footer(s, n)


def engagement_structure(prs, n=9):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Engagement Structure",
           "Clear ownership: one account manager, two embedded specialists")

    cx = Inches(5.0)

    # Client node
    client_box = rect(s, cx - Inches(1.6), Inches(1.85), Inches(3.2), Inches(0.70),
                      NAVY, line=BORDER)
    cb = client_box.text_frame
    cb.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = cb.paragraphs[0]
    cp.text = f"{CLIENT} — Project Owner"
    cp.font.size = Pt(11.5)
    cp.font.bold = True
    cp.font.color.rgb = WHITE
    cp.alignment = PP_ALIGN.CENTER

    # AM node
    am_box = rect(s, cx - Inches(1.6), Inches(3.00), Inches(3.2), Inches(0.70),
                  BLUE, line=BORDER)
    ab = am_box.text_frame
    ab.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = ab.paragraphs[0]
    ap.text = "Synvoric — Account Manager"
    ap.font.size = Pt(11.5)
    ap.font.bold = True
    ap.font.color.rgb = WHITE
    ap.alignment = PP_ALIGN.CENTER

    # Resource nodes
    res_y = Inches(4.20)
    for i, (name, role) in enumerate([("Ayush Shukla", "UX / UI Designer"),
                                       ("Sarthak Shukla", "Frontend Developer")]):
        x = cx - Inches(1.8 + (-i * 1.9))  # i=0 left, i=1 right
        if i == 0:
            x = cx - Inches(1.85)
        else:
            x = cx + Inches(0.05)
        box = rect(s, x, res_y, Inches(1.80), Inches(0.70), LIGHT_BG, line=BORDER)
        bt = box.text_frame
        bt.vertical_anchor = MSO_ANCHOR.MIDDLE
        bt.margin_left = Inches(0.10)
        np = bt.paragraphs[0]
        np.text = name
        np.font.size = Pt(10.5)
        np.font.bold = True
        np.font.color.rgb = NAVY
        np.alignment = PP_ALIGN.CENTER
        rp = bt.add_paragraph()
        rp.text = role
        rp.font.size = Pt(8.5)
        rp.font.color.rgb = SLATE
        rp.alignment = PP_ALIGN.CENTER

    # Connector lines (vertical between client & AM, then split to two)
    line_h(s, cx - Inches(0.01), Inches(2.55), Inches(0.02), color=SLATE, thickness=0.50)
    # using a thin vertical rectangle
    rect(s, cx - Inches(0.01), Inches(2.55), Inches(0.02), Inches(0.45), SLATE)
    # AM to resources horizontal line
    rect(s, cx - Inches(0.95), Inches(3.90), Inches(1.90), Inches(0.02), SLATE)
    # AM down stub
    rect(s, cx - Inches(0.01), Inches(3.70), Inches(0.02), Inches(0.20), SLATE)
    # Verticals down to each resource
    rect(s, cx - Inches(0.95), Inches(3.92), Inches(0.02), Inches(0.28), SLATE)
    rect(s, cx + Inches(0.95), Inches(3.92), Inches(0.02), Inches(0.28), SLATE)

    # Side annotations
    notes = [
        ("Weekly cadence", "Status report every Friday"),
        ("Single point of contact", "All escalations via Account Manager"),
        ("No team management overhead", "Synvoric handles payroll and oversight"),
    ]
    for i, (h, b) in enumerate(notes):
        y = Inches(1.85 + i * 1.05)
        txt(s, ML, y, Inches(2.0), Inches(0.25), h, 10, bold=True, color=NAVY)
        txt(s, ML, y + Inches(0.27), Inches(2.0), Inches(0.55), b, 9.5, color=SLATE)

    footer(s, n)


# ── candidate profiles ────────────────────────────────────────────────

def profile_slide(prs, n, name, role, action_title, meta: dict):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Candidate Profile", action_title)

    # Left card (navy)
    card_x, card_y = ML, BODY_Y
    card_w, card_h = Inches(2.45), Inches(3.30)
    rect(s, card_x, card_y, card_w, card_h, NAVY)

    # Photo or initials circle
    photo_path = PHOTOS.get(name)
    circle_x = card_x + (card_w - Inches(1.00)) / 2
    circle_y = card_y + Inches(0.25)
    if photo_path and photo_path.exists():
        s.shapes.add_picture(str(photo_path), circle_x, circle_y, Inches(1.00), Inches(1.00))
    else:
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, circle_x, circle_y, Inches(1.00), Inches(1.00))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.text = "".join(p[0] for p in name.split()[:2])
        cp.font.size = Pt(22)
        cp.font.bold = True
        cp.font.color.rgb = WHITE
        cp.alignment = PP_ALIGN.CENTER

    txt(s, card_x, card_y + Inches(1.40), card_w, Inches(0.32),
        name, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, card_x, card_y + Inches(1.72), card_w, Inches(0.28),
        role, 10, color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    line_h(s, card_x + Inches(0.95), card_y + Inches(2.05), Inches(0.55),
           color=BLUE, thickness=0.025)
    txt(s, card_x, card_y + Inches(2.20), card_w, Inches(0.25),
        meta["experience"], 10, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txt(s, card_x, card_y + Inches(2.50), card_w, Inches(0.25),
        meta["location"], 9, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER)
    if meta.get("link"):
        txt(s, card_x, card_y + Inches(2.78), card_w, Inches(0.25),
            meta["link"], 8, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.CENTER, italic=True)

    # Right column
    rx = card_x + card_w + Inches(0.35)
    rw = Inches(5.20)

    txt(s, rx, BODY_Y, rw, Inches(0.22), "SUMMARY", 8.5, bold=True, color=BLUE)
    txt(s, rx, BODY_Y + Inches(0.25), rw, Inches(0.85), meta["summary"], 10.5, color=BODY)

    txt(s, rx, BODY_Y + Inches(1.10), rw, Inches(0.22),
        "EXPERIENCE HIGHLIGHTS", 8.5, bold=True, color=BLUE)
    tf = txt(s, rx, BODY_Y + Inches(1.35), rw, Inches(1.30), "", 10)
    bullets(tf, meta["highlights"], size=10, spacing=5)

    txt(s, rx, BODY_Y + Inches(2.70), rw, Inches(0.22),
        "CORE SKILLS", 8.5, bold=True, color=BLUE)
    txt(s, rx, BODY_Y + Inches(2.95), rw, Inches(0.45), meta["skills"], 10, color=BODY)

    footer(s, n)


# ── commercials ───────────────────────────────────────────────────────

def pricing(prs, n=14):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Commercial Terms",
           "$7,000 USD / month for a full-time, dedicated 2-person team")

    # Table
    rows = [
        ["#", "Role", "Resource", "Allocation", "Monthly Rate (USD)"],
        ["1", "UX / UI Designer", "Ayush Shukla", "Full-time", f"$ {RATE_USD:,}"],
        ["2", "Frontend Developer", "Sarthak Shukla", "Full-time", f"$ {RATE_USD:,}"],
    ]
    row_heights = [Inches(0.40), Inches(0.50), Inches(0.50)]
    table_h = sum(row_heights, Inches(0))
    tbl = s.shapes.add_table(len(rows), 5, ML, BODY_Y, CW, table_h).table
    widths = [Inches(0.40), Inches(2.20), Inches(2.10), Inches(1.50), Inches(2.36)]
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = w
    for ri, h in enumerate(row_heights):
        tbl.rows[ri].height = h
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            style_cell(cell, fill=NAVY if ri == 0 else None)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Calibri"
                p.font.size = Pt(9.5 if ri == 0 else 10.5)
                p.font.bold = ri == 0
                p.font.color.rgb = WHITE if ri == 0 else BODY
                if ci in (0, 3, 4):
                    p.alignment = PP_ALIGN.CENTER

    # Total bar
    total_y = BODY_Y + table_h + Inches(0.10)
    bar = rect(s, ML, total_y, CW, Inches(0.55), NAVY)
    btf = bar.text_frame
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.margin_left = Inches(0.20)
    btf.margin_right = Inches(0.20)
    bp = btf.paragraphs[0]
    bp.text = "TOTAL MONTHLY INVESTMENT"
    bp.font.size = Pt(11)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.alignment = PP_ALIGN.LEFT
    # Right amount as separate overlay
    amt = s.shapes.add_textbox(ML, total_y, CW, Inches(0.55))
    atf = amt.text_frame
    atf.vertical_anchor = MSO_ANCHOR.MIDDLE
    atf.margin_right = Inches(0.20)
    ap = atf.paragraphs[0]
    ap.text = f"$ {MONTHLY_TOTAL:,}  /  month"
    ap.font.size = Pt(15)
    ap.font.bold = True
    ap.font.color.rgb = WHITE
    ap.alignment = PP_ALIGN.RIGHT

    # Payment & Terms grid
    terms_y = total_y + Inches(0.80)
    txt(s, ML, terms_y, CW, Inches(0.22), "PAYMENT & TERMS", 8.5, bold=True, color=BLUE)
    terms = [
        ("Billing", "Monthly in arrears, invoice due in 15 days"),
        ("Notice", "30-day rolling notice from either party"),
        ("Replacement", "Guaranteed within 30 days at no extra cost"),
        ("Compliance", "Synvoric manages payroll and contracts"),
    ]
    for i, (k, v) in enumerate(terms):
        col, row = i % 2, i // 2
        x = ML + Inches(col * 4.30)
        y = terms_y + Inches(0.28 + row * 0.36)
        txt(s, x, y, Inches(1.20), Inches(0.30), k, 10, bold=True, color=NAVY)
        txt(s, x + Inches(1.20), y, Inches(2.95), Inches(0.30), v, 10, color=BODY)

    footer(s, n,
           source=f"Annual run-rate: $ {MONTHLY_TOTAL * 12:,} USD  ·  all fees exclusive of applicable taxes")


def engagement_model(prs, n=15):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Engagement Model",
           "Four-phase engagement from kickoff to steady-state delivery")
    steps = [
        ("Week 1", "Kickoff", "Agreement signed, access provisioned, sprint cadence aligned"),
        ("Week 2", "Onboarding", "Resources embedded in your tools, processes and rituals"),
        ("Ongoing", "Delivery", "Designs by Ayush, implementation by Sarthak; weekly demos"),
        ("Quarterly", "Review & Scale", "Performance review with option to scale team up or down"),
    ]
    for i, (phase, title, desc) in enumerate(steps):
        x = ML + Inches(i * 2.15)
        box = rect(s, x, BODY_Y, Inches(1.95), Inches(2.25), LIGHT_BG, line=BORDER)
        line_h(s, x, BODY_Y, Inches(1.95), color=BLUE, thickness=0.04)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.12), Inches(1.70), Inches(0.25),
            phase.upper(), 9, bold=True, color=BLUE)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.40), Inches(1.70), Inches(0.32),
            title, 12, bold=True, color=NAVY)
        txt(s, x + Inches(0.12), BODY_Y + Inches(0.80), Inches(1.70), Inches(1.30),
            desc, 10, color=BODY)
    footer(s, n)


def risk_mitigation(prs, n=16):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Risk Mitigation",
           "Built-in safeguards: replacement, IP protection, and continuity")
    risks = [
        ("Resource fit",
         "If a developer does not meet expectations, we provide a free replacement within 30 days.",
         "Replacement Guarantee"),
        ("IP protection",
         "All deliverables, code, and designs are owned by Consumables Solutions upon payment.",
         "IP Assignment in MSA"),
        ("Confidentiality",
         "Both contractors sign NDAs covering Consumables Solutions information.",
         "NDA-Backed"),
        ("Business continuity",
         "Synvoric retains backup talent in our network for immediate switchover if needed.",
         "Bench-Backed"),
        ("Compliance",
         "Synvoric handles all payroll, tax, and statutory obligations — zero burden on you.",
         "Fully Managed"),
        ("Performance",
         "Weekly status reports and monthly review calls ensure issues are caught early.",
         "Transparent Reporting"),
    ]
    for i, (cat, body, tag) in enumerate(risks):
        col, row = i % 3, i // 3
        x = ML + Inches(col * 2.86)
        y = BODY_Y + Inches(row * 1.50)
        box = rect(s, x, y, Inches(2.70), Inches(1.35), LIGHT_BG, line=BORDER)
        line_h(s, x, y, Inches(2.70), color=BLUE, thickness=0.04)
        txt(s, x + Inches(0.12), y + Inches(0.10), Inches(2.45), Inches(0.30),
            cat, 11, bold=True, color=NAVY)
        txt(s, x + Inches(0.12), y + Inches(0.42), Inches(2.45), Inches(0.65),
            body, 9.5, color=BODY)
        tag_box = rect(s, x + Inches(0.12), y + Inches(1.05), Inches(1.6), Inches(0.20),
                       BLUE)
        ttf = tag_box.text_frame
        ttf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ttf.margin_left = Inches(0.08)
        tp = ttf.paragraphs[0]
        tp.text = tag.upper()
        tp.font.size = Pt(7.5)
        tp.font.bold = True
        tp.font.color.rgb = WHITE
    footer(s, n)


def next_steps(prs, n=17):
    s = prs.slides.add_slide(blank_layout(prs))
    header(s, "Next Steps",
           "Five steps to a deployed team, completed within one week")
    steps = [
        ("Day 0–1", "Review proposal and candidate profiles with your team"),
        ("Day 1–2", "Confirm role requirements and target start date with Synvoric"),
        ("Day 3–4", "Execute MSA and Statement of Work"),
        ("Day 4–5", "Resources onboarded; access, tools and standups confirmed"),
        ("Day 5+", "Delivery begins — weekly status reports start immediately"),
    ]
    for i, (when, what) in enumerate(steps):
        y = BODY_Y + Inches(i * 0.50)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, ML, y, Inches(0.40), Inches(0.40))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.text = str(i + 1)
        cp.font.size = Pt(11)
        cp.font.bold = True
        cp.font.color.rgb = WHITE
        cp.alignment = PP_ALIGN.CENTER

        txt(s, Inches(1.30), y + Inches(0.04), Inches(1.50), Inches(0.32),
            when, 10, bold=True, color=BLUE)
        txt(s, Inches(2.90), y + Inches(0.04), Inches(5.65), Inches(0.32),
            what, 11, color=BODY)

    cta = rect(s, ML, Inches(4.65), CW, Inches(0.45), NAVY)
    ctf = cta.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.20)
    cp = ctf.paragraphs[0]
    cp.text = "Contact: anshul.agrawal@synvoric.com  ·  www.synvoric.com  ·  Proposal valid for 30 days"
    cp.font.size = Pt(11)
    cp.font.color.rgb = WHITE
    footer(s, n)


def closing(prs, n=18):
    s = prs.slides.add_slide(blank_layout(prs))
    set_bg(s, NAVY)
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(3.05), Inches(1.15), width=Inches(3.9))
    txt(s, Inches(1), Inches(2.80), Inches(8), Inches(0.5),
        "Thank You", 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.40), Inches(8), Inches(0.35),
        f"We look forward to partnering with {CLIENT}", 13,
        color=RGBColor(0xCB, 0xD5, 0xE1), align=PP_ALIGN.CENTER)
    txt(s, Inches(1), Inches(3.95), Inches(8), Inches(0.30),
        "anshul.agrawal@synvoric.com  ·  synvoric.com", 11, color=BLUE, align=PP_ALIGN.CENTER)
    footer(s, n, dark=True)


# ── build ─────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    cover(prs)
    executive_summary(prs, 2)
    agenda(prs, 3)
    about_synvoric(prs, 4)
    why_synvoric(prs, 5)
    divider(prs, "01", "The Proposal",
            f"Understanding {CLIENT}'s need and our recommended approach", 6)
    client_context(prs, 7)
    proposed_team(prs, 8)
    engagement_structure(prs, 9)
    divider(prs, "02", "Candidate Profiles",
            "Meet the dedicated 2-person team", 10)
    profile_slide(prs, 11, "Ayush Shukla", "UX / UI Designer",
                  "Ayush Shukla — 2+ years designing conversion-focused web and mobile products",
                  {
                      "location": "Meerut, India",
                      "experience": "2+ Years Experience",
                      "link": "ayushshukla18.framer.website",
                      "summary": (
                          "UX/UI Designer with 2+ years designing web and mobile products in Figma. "
                          "Experienced across recruitment, healthcare, wellness, and SaaS domains. "
                          "Strong in user flows, wireframes, and high-fidelity responsive design."
                      ),
                      "highlights": [
                          "SeekHelpers — 100+ responsive screens for hiring platform (2025–2026)",
                          "fyp.life — Mental wellness website UX, accessibility-focused",
                          "Vlippr — 150+ mobile screens, voice-cloning app interaction design",
                          "Cross-functional collaboration with PMs and engineering teams",
                      ],
                      "skills": "Figma · Adobe XD · Wireframing · Prototyping · Design Systems · User Research · HTML/CSS",
                  })
    profile_slide(prs, 12, "Sarthak Shukla", "Frontend Developer",
                  "Sarthak Shukla — 4 years shipping React and React Native at scale",
                  {
                      "location": "Meerut, India",
                      "experience": "4 Years Experience",
                      "link": "github.com/sarthakshukla99",
                      "summary": (
                          "Frontend Developer with 4 years building scalable web and mobile apps. "
                          "Expert in translating Figma designs to production-ready React and "
                          "React Native code. Performer of the Year (2023–24) at Shilsha Technologies."
                      ),
                      "highlights": [
                          "Shilsha Technologies — AI/ML analytics platform & e-commerce app",
                          "Built seller dashboard: inventory, orders, analytics in React.js",
                          "Fodrix — Reusable JS library, MFA; UX rating raised 55% → 85%",
                          "GenAI PDF Assistant using OpenAI API and vector search",
                      ],
                      "skills": "React.js · Next.js · TypeScript · React Native · Redux · Tailwind CSS · REST APIs · Git",
                  })
    divider(prs, "03", "Commercial Terms",
            "Transparent pricing, flexible engagement, built-in safeguards", 13)
    pricing(prs, 14)
    engagement_model(prs, 15)
    risk_mitigation(prs, 16)
    next_steps(prs, 17)
    closing(prs, 18)

    prs.save(str(OUT))
    print(f"Created {OUT} ({TOTAL_SLIDES} slides)")


if __name__ == "__main__":
    ASSETS.mkdir(parents=True, exist_ok=True)
    build()
